#!/usr/bin/env python3
"""Build a plain, academic-style slide deck (.pptx).

Design: white background, navy headings, a single muted-gold accent rule,
serif titles, sans body, generous whitespace. No animation, no clutter --
the look of a conference / lab talk built in Beamer.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- palette ----
NAVY   = RGBColor(0x1F, 0x3A, 0x5F)
GOLD   = RGBColor(0xB5, 0x79, 0x1F)
INK    = RGBColor(0x22, 0x28, 0x2E)
GREY   = RGBColor(0x5B, 0x66, 0x70)
LIGHT  = RGBColor(0xEA, 0xEC, 0xEF)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

TITLE_FONT = "Georgia"
BODY_FONT  = "Calibri"
MONO_FONT  = "Consolas"

EMU = 914400
SW, SH = Inches(13.333), Inches(7.5)  # 16:9

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]

FOOTER = "Agent Skills for ML Engineering — PengTao Lab, UCSD"


def add_slide():
    return prs.slides.add_slide(BLANK)


def fill_bg(slide, color=WHITE):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def box(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf


def style_run(r, size, color=INK, bold=False, font=BODY_FONT, italic=False):
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font


def rule(slide, l, t, w, color=GOLD, h=Pt(3)):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def eyebrow(slide, text, l=Inches(0.9), t=Inches(0.62)):
    tb, tf = box(slide, l, t, Inches(11.5), Inches(0.4))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text.upper()
    style_run(r, 12.5, GOLD, bold=True, font=BODY_FONT)
    r.font._rPr.set("spc", "200")  # letter spacing
    return tb


def heading(slide, text, t=Inches(1.0), size=30):
    tb, tf = box(slide, Inches(0.9), t, Inches(11.5), Inches(1.1))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    style_run(r, size, NAVY, bold=True, font=TITLE_FONT)
    rule(slide, Inches(0.92), t + Inches(0.92), Inches(0.9))
    return tb


_SLIDE_NO = [1]  # slide 1 is the title (no footer); footer() auto-numbers from 2


def footer(slide, idx=None):
    _SLIDE_NO[0] += 1
    n = _SLIDE_NO[0]
    tb, tf = box(slide, Inches(0.9), Inches(7.02), Inches(9.5), Inches(0.35))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = FOOTER
    style_run(r, 9, GREY, font=BODY_FONT)
    nb, ntf = box(slide, Inches(12.0), Inches(7.02), Inches(0.9), Inches(0.35))
    np_ = ntf.paragraphs[0]; np_.alignment = PP_ALIGN.RIGHT
    nr = np_.add_run(); nr.text = str(n)
    style_run(nr, 10, GREY, font=BODY_FONT)


def bullets(slide, items, l=Inches(0.95), t=Inches(2.15), w=Inches(11.4),
            h=Inches(4.4), size=18, gap=10):
    tb, tf = box(slide, l, t, w, h)
    first = True
    for it in items:
        lvl = it[0] if isinstance(it, tuple) else 0
        text = it[1] if isinstance(it, tuple) else it
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lvl
        p.space_after = Pt(gap)
        p.space_before = Pt(0)
        marker = "—  " if lvl == 0 else "·  "
        mr = p.add_run(); mr.text = marker
        style_run(mr, size, GOLD if lvl == 0 else GREY, bold=(lvl == 0))
        # support **bold** segments
        segs = text.split("**")
        for i, seg in enumerate(segs):
            if seg == "":
                continue
            r = p.add_run(); r.text = seg
            style_run(r, size - (lvl * 1), INK if lvl == 0 else GREY,
                      bold=(i % 2 == 1))
    return tb


# ----------------------------------------------------------------------
# SLIDE 1 -- title
# ----------------------------------------------------------------------
s = add_slide(); fill_bg(s)
rule(s, Inches(0.9), Inches(2.05), Inches(1.6), color=GOLD, h=Pt(4))
tb, tf = box(s, Inches(0.9), Inches(2.25), Inches(11.5), Inches(2.2))
p = tf.paragraphs[0]
r = p.add_run(); r.text = "Do Skills Make ML-Engineering Agents"
style_run(r, 40, NAVY, bold=True, font=TITLE_FONT)
p2 = tf.add_paragraph()
r = p2.add_run(); r.text = "Measurably Better?"
style_run(r, 40, NAVY, bold=True, font=TITLE_FONT)
tb2, tf2 = box(s, Inches(0.9), Inches(4.35), Inches(11.5), Inches(0.6))
p = tf2.paragraphs[0]
r = p.add_run(); r.text = "A system for building, discovering, and evaluating Agent Skills"
style_run(r, 19, GREY, italic=True, font=TITLE_FONT)
tb3, tf3 = box(s, Inches(0.9), Inches(5.55), Inches(11.5), Inches(1.2))
p = tf3.paragraphs[0]
r = p.add_run(); r.text = "Kuntal Kokate"
style_run(r, 16, INK, bold=True)
p = tf3.add_paragraph()
r = p.add_run(); r.text = "PengTao Lab, University of California, San Diego"
style_run(r, 14, GREY)
p = tf3.add_paragraph()
r = p.add_run(); r.text = "June 2026"
style_run(r, 13, GREY)

# ----------------------------------------------------------------------
# SLIDE 2 -- outline
# ----------------------------------------------------------------------
s = add_slide(); fill_bg(s)
eyebrow(s, "Outline")
heading(s, "Talk roadmap")
bullets(s, [
    "**The question** — do skills change what an agent builds, not just what it says?",
    "**The north star** — a library of skills the agent draws from at task time",
    "**Part I · The agents** — building, discovering, and testing skills",
    "**Part II · The evaluation** — a two-stage pipeline (Stage 1 local, Stage 2 A/B)",
    "**The MLEvolve A/B framework** — skill injection, trustworthy metric, runtime",
    "**Results & what they reveal** — SAMSum, gsm8k, and gaps in the skill builder",
    "**Conclusion** and the full sweep ahead",
], size=19, gap=14)
footer(s, 2)

# ----------------------------------------------------------------------
# SLIDE 3 -- motivation
# ----------------------------------------------------------------------
s = add_slide(); fill_bg(s)
eyebrow(s, "Motivation")
heading(s, "Skills are easy to assert, hard to verify")
bullets(s, [
    "An **Agent Skill** is a Markdown playbook injected into an agent's context, using progressive disclosure: short description always loaded, body on trigger, references as needed.",
    "Skills are cheap to write and easy to share — but their effect is usually **asserted, not measured**.",
    "The operational question of this project:",
    (1, "Does an MLE agent build a **measurably better** pipeline *with* a skill than without it?"),
    (1, "And if so, **where** in the pipeline does the help land?"),
    "Answering it needs three capabilities: **build** good skills, **discover** existing ones, and **evaluate** their downstream effect.",
], size=18, gap=12)
footer(s, 3)

# ----------------------------------------------------------------------
# SLIDE -- north star / vision
# ----------------------------------------------------------------------
s = add_slide(); fill_bg(s)
eyebrow(s, "The north star")
heading(s, "A skill library the agent draws from")
bullets(s, [
    "The end state: the **ai-skill-builder** agent produces many skills into one **shared library**.",
    "At task time the **whole library** is handed to the MLE agent — not a single hand-picked skill.",
    "While solving, the agent **selects a few skills of its own choosing**, per step, and loads only what it needs.",
    "Progressive disclosure at library scale: **Discovery** (see all) → **Activation** (pick skills) → **Execution** (pick references).",
    "It closes the loop — **build → library → select → measure** — and the evaluation exists to prove the library helps.",
], size=18, gap=13)
footer(s)

# ----------------------------------------------------------------------
# SLIDE 4 -- three agents
# ----------------------------------------------------------------------
s = add_slide(); fill_bg(s)
eyebrow(s, "Part I · The agents")
heading(s, "Three OpenClaw agents")
# simple 3-row table feel via bullets with sub-detail
bullets(s, [
    "**ai-skill-builder** — synthesises a SKILL.md (+ references, scripts) from one documentation URL.",
    (1, "Phase 1.5 · v1.4.0 · the LLM writes the body; Python assembles the frontmatter."),
    "**ai-skill-scout** — searches GitHub for existing skills, ranks by trust, installs after a scan.",
    (1, "Phase 1 complete · nothing reaches the workspace unscanned."),
    "**skill-tester** — drives Stage-1 local evaluation; an MCP sidecar captures every prompt.",
    (1, "In production · measures triggering and functional lift vs a baseline agent."),
], size=18, gap=10)
footer(s, 4)

# ----------------------------------------------------------------------
# SLIDE 5 -- builder
# ----------------------------------------------------------------------
s = add_slide(); fill_bg(s)
eyebrow(s, "Part I · Skill-builder")
heading(s, "One URL in, one skill out")
bullets(s, [
    "**Pipeline:** resolve → fetch → extract hints → IPI scan → plan → synthesise → triggering eval → assemble → validate → write.",
    "**The model never writes the frontmatter** — name, install, MCP declarations, and a provenance hash are assembled deterministically in Python.",
    "**Sources, ranked for safety:** docs + README + examples always; Stack Exchange and closed bug issues opt-in; unmoderated **forums excluded** (prompt-injection risk).",
    "**Gates:** a 60-pattern security scanner (7 threat categories), an IPI scan, a line cap, and a shell-command fabrication check — all pre-write.",
    "**Result:** the peft-tuning skill scored a triggering **F1 = 1.000** over 60 judge calls.",
], size=18, gap=12)
footer(s, 5)

# ----------------------------------------------------------------------
# SLIDE 6 -- scout + tester
# ----------------------------------------------------------------------
s = add_slide(); fill_bg(s)
eyebrow(s, "Part I · Discover & test")
heading(s, "Find a skill safely; check that it fires")
bullets(s, [
    "**Skill-scout** — GitHub-only search, LLM query expansion, then rank by (trust, class, completeness, stars).",
    (1, "Quarantine to /tmp → the **same** 60-pattern scan → user approval → install. One source of truth for danger."),
    "**Skill-tester / Stage 1** — a fast, CI-style check on every build (~5 min, ~$1):",
    (1, "Triggering F1 (target ≥ 0.85) · functional pass-rate · lift · citation rate (≥ 80%) · cost ratio (≤ 2×)."),
    "**Reliability bar:** the **description is the #1 lever** — wording alone swings invocation > 10× with capability fixed.",
], size=18, gap=12)
footer(s, 6)

# ----------------------------------------------------------------------
# SLIDE 7 -- two-stage
# ----------------------------------------------------------------------
s = add_slide(); fill_bg(s)
eyebrow(s, "Part II · The evaluation")
heading(s, "Two stages: does it fire — does it help?")
bullets(s, [
    "**Stage 1 · local (isolation).** Does the skill fire and surface the right facts?",
    (1, "No agent loop · triggering + functional pass-rate · runs on every build · ~5 min, ~$1."),
    (1, "Locked, in production. Rejects weak skills before we spend GPU hours."),
    "**Stage 2 · A/B (full run).** Does the agent build a better pipeline, and where?",
    (1, "Full agent run, paired with-skill / without-skill · 3 seeds × 3 tasks = 18 trajectories."),
    (1, "Held-out grader + per-stage attribution · a pre-ship gate · harness validated end-to-end."),
], size=18, gap=11)
footer(s, 7)

# ----------------------------------------------------------------------
# SLIDE 8 -- MLEvolve setup
# ----------------------------------------------------------------------
s = add_slide(); fill_bg(s)
eyebrow(s, "Part II · Stage 2 · MLEvolve")
heading(s, "A paired with / without-skill A/B")
bullets(s, [
    "**Agent:** MLEvolve-generic — a Monte-Carlo graph search over candidate code nodes, driven by deepseek-v4-pro.",
    (1, "Subprocess-per-node design avoids the fork-after-CUDA crash that retired the earlier AIDE agent."),
    "**Held fixed across cells:** same backbone (Qwen2.5-3B-Instruct), same task, same seed, same skill library.",
    (1, "The *only* difference is whether the library is available; the agent selects from it (empty library = baseline)."),
    (1, "Today's spike held one skill (peft-tuning) fixed; the goal is the full library the agent chooses from."),
    "**Tasks:** samsum (ROUGE-L) · gsm8k (exact-match) · boolq (accuracy).",
    "**Measured:** L1 outcome (held-out grader + Lift) · L2 *where* it helped · L3 cost · **selection quality**.",
], size=18, gap=10)
footer(s, 8)

# ----------------------------------------------------------------------
# SLIDE 9 -- skill injection
# ----------------------------------------------------------------------
s = add_slide(); fill_bg(s)
eyebrow(s, "Part II · Skill injection")
heading(s, "Library in → the agent picks a few")
bullets(s, [
    "MLEvolve does single-shot code generation — it cannot open a file to read a skill, so the injector hands it the **whole library** and lets it choose.",
    "**Tier 0 · Discovery (every node).** A ~150-token catalogue of the *whole library* — names, one-liners, reference files. The agent sees everything available.",
    "**Tier 1 · Activation (per node).** A temperature-0 selector picks **which skills** this step needs — a few of the agent's choice, not all of them.",
    "**Tier 2 · Execution.** It loads only the relevant references — never the whole library into every node.",
    "Every selection is logged, so we also measure **whether the agent picked the right skills** — selection precision / recall.",
], size=18, gap=12)
footer(s, 9)

# ----------------------------------------------------------------------
# SLIDE 10 -- trustworthy metric
# ----------------------------------------------------------------------
s = add_slide(); fill_bg(s)
eyebrow(s, "Part II · The harness")
heading(s, "A number you can trust")
bullets(s, [
    "An agent that reports its own score will learn to report a good one — so the headline never comes from the agent.",
    "**Independent held-out grader:** after exit, mleval.grader recomputes the metric from the preserved submission.csv against references the agent never saw.",
    "The self-reported validation score is kept only as the **tree-search signal** and a drift diagnostic (one 0.81 self-report graded out as invalid).",
    "**Runtime:** UCSD Nautilus NRP Kubernetes — one image per agent, one Pod per (task × cell × seed).",
    "**Sidecar patches:** seed · prompt-logger · token-budget (16k→32k) · skill-injector (the only one carrying the treatment); a build-time de-Kaggle patch stops off-task drift.",
], size=18, gap=11)
footer(s, 10)

# ----------------------------------------------------------------------
# SLIDE 11 -- results
# ----------------------------------------------------------------------
s = add_slide(); fill_bg(s)
eyebrow(s, "Preliminary results · spike-018")
heading(s, "SAMSum: what the grader actually shows")
bullets(s, [
    "**spike-018** ran the SAMSum 2×2 — with / without skill × seeds 0–1 — scored by the independent held-out grader over all 819 test ids.",
    "**Only 2 of 4 cells produced a valid submission:** an output-contract bug (wrong columns) sank one cell in *each* arm — so **no seed-matched pair, no Lift yet**.",
    "**Valid held-out ROUGE-L:** with-skill (seed 1) **0.288** · without-skill (seed 0) **0.227** — cross-seed, indicative only.",
    "**The grader earned its keep:** it flagged two inflated self-reports (0.94, 0.81) as **invalid** — exactly the off-task drift it exists to catch.",
    "**Where the skill shows:** self-correction 0.6 vs 0.15 · far less redundant work (3–4 vs 9–12 nodes) · significant stage shift (χ²=65, p≈3e-9).",
], size=18, gap=11)
footer(s, 11)

# ----------------------------------------------------------------------
# SLIDE -- gsm8k
# ----------------------------------------------------------------------
s = add_slide(); fill_bg(s)
eyebrow(s, "Preliminary results · gsm8k")
heading(s, "gsm8k: pipeline proven, A/B still blocked")
bullets(s, [
    "Rewrote gsm8k to an **MLE-Bench-aligned held-out re-split** — labels withheld on test, a carved val for the agent's own signal.",
    "**spike-023 validated the held-out pipeline end-to-end:** a without-skill baseline of **0.61 exact-match (809/1319)**, graded from the preserved 1319-row submission (not auto-persisted).",
    "**Caught + fixed a skill-router regression:** a 3 KB harness-rules header pushed the task past the selector's 1500-char cap → it picked **no skills** → the treatment arm was effectively empty. Fixed (sentinel + cap → 6000).",
    "**Bottleneck is generation throughput, not training:** every node trains in ~10–15 min; timeouts hit the 1319-example decode. One node ratcheted batch 4→32 / max_new_tokens 512→200 to fit 52 min.",
    "**No clean paired A/B yet** — earlier runs were killed by the 60-min per-exec cap (we mark unfinished nodes buggy). Next: raise the cap and rerun.",
], size=18, gap=10)
footer(s)

# ----------------------------------------------------------------------
# SLIDE -- skill-builder issues (from vllm)
# ----------------------------------------------------------------------
s = add_slide(); fill_bg(s)
eyebrow(s, "What gsm8k revealed · the skill builder")
heading(s, "A skill can fire and still not help")
bullets(s, [
    "**vLLM was never the right tool for gsm8k** — the task prescribes HF AutoModelForCausalLM + batched generate, and the agent correctly used HF in *both* arms. 'Agent ignored vLLM' was a mis-read, not a skill defect.",
    "**Genre self-exclusion:** vllm-inference's *NOT-for* section says 'SFT / LoRA training → use transformers' — so on a LoRA task the skill **scoped itself out** and steered the agent away. Counterproductive, not just inert.",
    "**It passed triggering anyway** (F1 ≥ 0.85, selected 10/11) — the judge competes vs 5 canned decoys, not the *real* siblings, so 'fine-tune+eval → peft-tuning' was never tested.",
    "**The decisive failure was functional, not selection:** both arms ran HF; runs died on HF hygiene (use_cache off under grad-checkpointing, loose max_new_tokens, over-eval) that peft-tuning (stale 1.3.0) doesn't teach.",
    "**Creator fixes:** (1) a **library-aware** triggering eval (judge vs real siblings) and (2) a **functional** eval that rewards operational gotchas — not just tightening descriptions.",
], size=17, gap=9)
footer(s)

# ----------------------------------------------------------------------
# SLIDE 12 -- conclusion
# ----------------------------------------------------------------------
s = add_slide(); fill_bg(s)
eyebrow(s, "Conclusion")
heading(s, "The loop is closed — now we run it")
bullets(s, [
    "**Delivered:** an end-to-end loop for Agent Skills — build, discover, evaluate — with the A/B framework as the research contribution.",
    "**Next:** raise the per-exec cap and run the full paired sweep (samsum · gsm8k · boolq) — mean Lift with a 95% CI, L2 stage attribution, and skill-selection precision/recall.",
    "**Harden the builder too:** library-aware triggering + a functional choreography eval, so 'fires' implies 'helps'.",
    "**Takeaways:** (1) a trustworthy metric needs a grader the agent can't reach — it already caught real drift; (2) discovery is necessary but not sufficient — a fired skill can still be inert; (3) whether a builder-made library + agent selection helps is exactly what the fixed harness now answers.",
], size=17, gap=11)
footer(s)

out = "project-slides.pptx"
prs.save(out)
print("wrote", out, "with", len(prs.slides._sldIdLst), "slides")
